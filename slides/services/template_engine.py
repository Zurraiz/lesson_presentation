import os
import requests
from io import BytesIO
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from django.conf import settings

class TemplateManager:
    """
    Manages PowerPoint template operations: inspection and generation.
    """
    
    def __init__(self, template_path):
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template not found: {template_path}")
        self.template_path = template_path
        self.prs = Presentation(template_path)

    def analyze_template(self):
        """
        Returns a JSON-serializable structure defining available layouts and their placeholders.
        """
        layouts_info = []
        
        for idx, layout in enumerate(self.prs.slide_layouts):
            placeholders = []
            for ph in layout.placeholders:
                ph_type = ph.placeholder_format.type
                ph_name = ph.name
                ph_idx = ph.placeholder_format.idx
                
                is_generic = ph_type in [PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.BODY]
                
                placeholders.append({
                    "index": ph_idx,
                    "name": ph_name,
                    "type": str(ph_type), 
                    "is_title": "title" in ph_name.lower() or ph_type == PP_PLACEHOLDER.TITLE,
                    "is_image": is_generic or "picture" in ph_name.lower() or ph_type in [PP_PLACEHOLDER.PICTURE, PP_PLACEHOLDER.BITMAP],
                    "is_table": is_generic or "table" in ph_name.lower() or ph_type == PP_PLACEHOLDER.TABLE,
                    "is_chart": is_generic or "chart" in ph_name.lower() or ph_type == PP_PLACEHOLDER.CHART
                })
            
            layouts_info.append({
                "id": idx,
                "name": layout.name,
                "placeholders": placeholders
            })
            
        return {"layouts": layouts_info}

    def create_presentation(self, slides_data, output_filename=None):
        """
        Creates a new presentation based on the template and provided data.
        
        slides_data: List of dictionaries. Each dict should have:
            - layout_id: int (index of the layout to use)
            - content: dict mapping placeholder_index (str/int) to value
            
        Returns: Path to the saved file.
        """
        # Create a fresh presentation instance from the template
        new_prs = Presentation(self.template_path)
        
        # Remove any existing slides (templates sometimes have dummy slides)
        # Note: python-pptx doesn't support deleting common slides easily without low-level XML, 
        # but usually we just append. If the template file comes with slides, they remain.
        # Ideally, the template file should be empty of slides, just masters/layouts.
        
        for slide_data in slides_data:
            layout_id = slide_data.get("layout_id")
            content = slide_data.get("content", {})
            
            if layout_id is None:
                continue
                
            try:
                slide_layout = new_prs.slide_layouts[layout_id]
            except IndexError:
                # Fallback or skip
                continue
                
            slide = new_prs.slides.add_slide(slide_layout)
            self._fill_slide(slide, content)
            
        if not output_filename:
            output_filename = "generated_lesson.pptx"
            
        output_path = os.path.join(settings.MEDIA_ROOT, output_filename)
        new_prs.save(output_path)
        
        return output_path

    def _fill_slide(self, slide, content_map):
        """
        Fills a single slide's placeholders with content.
        Includes robust fallback logic (fuzzy matching) if AI provides wrong indices.
        """
        # 1. Map available placeholders by index and by type
        available_phs = {} # index -> shape
        by_type = {
            "title": [],
            "image": [],
            "table": [],
            "chart": [],
            "text": []
        }
        
        for ph in slide.placeholders:
            ph_idx = ph.placeholder_format.idx
            ph_type = ph.placeholder_format.type
            ph_name = ph.name.lower()
            available_phs[ph_idx] = ph
            
            # Categorize
            if ph_type == PP_PLACEHOLDER.TITLE or "title" in ph_name:
                by_type["title"].append(ph)
            elif ph_type in [PP_PLACEHOLDER.PICTURE, PP_PLACEHOLDER.BITMAP] or "picture" in ph_name:
                by_type["image"].append(ph)
            elif ph_type == PP_PLACEHOLDER.TABLE or "table" in ph_name:
                by_type["table"].append(ph)
            elif ph_type == PP_PLACEHOLDER.CHART or "chart" in ph_name:
                by_type["chart"].append(ph)
            else:
                by_type["text"].append(ph)

        # 2. Process content
        unmatched_content = []
        
        for raw_idx, value in content_map.items():
            try:
                look_idx = int(raw_idx)
            except:
                unmatched_content.append(value)
                continue
                
            # Attempt exact match
            if look_idx in available_phs:
                self._fill_shape(available_phs[look_idx], value)
                available_phs.pop(look_idx) # Mark as used
            else:
                unmatched_content.append(value)

        # 3. Fallback: Fuzzy match unmatched content to remaining placeholders
        for value in unmatched_content:
            target_type = "text"
            if isinstance(value, dict):
                v_type = value.get("type")
                if v_type in ["image", "table", "chart"]:
                    target_type = v_type
            
            # Priority 1: Direct type match (if available)
            if by_type[target_type]:
                ph = by_type[target_type].pop(0)
                self._fill_shape(ph, value)
            # Priority 2: Try "text" placeholder if it's a string
            elif target_type == "text" and by_type["text"]:
                ph = by_type["text"].pop(0)
                self._fill_shape(ph, value)

    def _fill_shape(self, shape, value):
        """Helper to route value to the correct insertion method for a shape."""
        if isinstance(value, str):
            # Text routing
            if not shape.has_text_frame:
                return
            shape.text = value
            self._apply_font_scaling(shape, value)
            
        elif isinstance(value, dict) and value.get("type") == "image":
            self._insert_image(shape, value.get("url"))
            
        elif isinstance(value, dict) and value.get("type") == "table":
            self._insert_table(shape, value)
            
        elif isinstance(value, dict) and value.get("type") == "chart":
            self._insert_chart(shape, value)

    def _apply_font_scaling(self, shape, text):
        """Dynamic Font Scaling logic."""
        try:
            text_len = len(text)
            if shape.text_frame.paragraphs:
                p = shape.text_frame.paragraphs[0]
                if not p.runs: p.add_run()
                
                if text_len > 300: font_size = Pt(12)
                elif text_len > 200: font_size = Pt(14)
                elif text_len > 100: font_size = Pt(18)
                else: font_size = Pt(24)
                
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = font_size
                        
            shape.text_frame.word_wrap = True
        except Exception as e:
            print(f"Error sizing text: {e}")

    def _insert_image(self, shape, image_url):
        """
        Inserts an image into a picture placeholder, preserving aspect ratio if possible.
        """
        # Download image to memory
        try:
            # Add User-Agent to avoid 403 Forbidden from sites like Wikimedia
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
            }
            response = requests.get(image_url, headers=headers, timeout=10)
            response.raise_for_status()
            
            # Use Pillow to ensure image is in a supported format (e.g. convert WEBP to PNG)
            image_stream = BytesIO(response.content)
            try:
                img = Image.open(image_stream)
                # Convert to RGB/RGBA then to PNG if not already a standard format
                # This fixes the "WEBP" not supported issue in many PPT versions
                converted_stream = BytesIO()
                img.save(converted_stream, format='PNG')
                converted_stream.seek(0)
                final_stream = converted_stream
            except Exception as img_err:
                print(f"Pillow conversion failed, trying raw stream: {img_err}")
                image_stream.seek(0)
                final_stream = image_stream

            # If it's a proper placeholder, insert_picture creates a NEW shape and replaces the placeholder
            if hasattr(shape, 'insert_picture'):
                shape.insert_picture(final_stream)
            else:
                # Fallback if it's not a picture placeholder but we want to force it? 
                # For now stick to strict placeholder behavior.
                pass
                
        except Exception as e:
            print(f"Failed to load image: {e}")
            # Optionally set text to error message
            if shape.has_text_frame:
                shape.text = f"[Error loading image]"

    def _insert_table(self, shape, data):
        """
        Inserts a table into a placeholder.
        data: {"headers": ["Col1", "Col2"], "rows": [["v1", "v2"], ...]}
        """
        try:
            headers = data.get("headers", [])
            rows = data.get("rows", [])
            
            num_rows = len(rows) + (1 if headers else 0)
            num_cols = len(headers) if headers else (len(rows[0]) if rows else 0)
            
            if num_rows == 0 or num_cols == 0:
                return
                
            # insert_table returns a GraphicFrame object containing the Table
            if hasattr(shape, 'insert_table'):
                table_frame = shape.insert_table(rows=num_rows, cols=num_cols)
                table = table_frame.table
                
                # Fill headers
                if headers:
                    for i, h in enumerate(headers):
                        table.cell(0, i).text = str(h)
                
                # Fill rows
                start_row = 1 if headers else 0
                for r_idx, row_data in enumerate(rows):
                    for c_idx, val in enumerate(row_data):
                        if c_idx < num_cols:
                            table.cell(r_idx + start_row, c_idx).text = str(val)
            else:
                print(f"Shape {shape.name} does not support table insertion.")
                
        except Exception as e:
            print(f"Failed to insert table: {e}")

    def _insert_chart(self, shape, data):
        """
        Inserts a chart into a placeholder.
        data: {"chart_type": "...", "categories": ["X", "Y"], "series": [{"name": "S1", "values": [10, 20]}]}
        """
        try:
            chart_data = CategoryChartData()
            chart_data.categories = data.get("categories", [])
            
            for s in data.get("series", []):
                chart_data.add_series(s.get("name", ""), s.get("values", []))
                
            # Default to Bar if not specified/invalid
            c_type_str = data.get("chart_type", "BAR_CLUSTERED").upper()
            c_type = getattr(XL_CHART_TYPE, c_type_str, XL_CHART_TYPE.BAR_CLUSTERED)
            
            if hasattr(shape, 'insert_chart'):
                shape.insert_chart(c_type, chart_data)
            else:
                print(f"Shape {shape.name} does not support chart insertion.")
                
        except Exception as e:
            print(f"Failed to insert chart: {e}")
