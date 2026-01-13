from pptx import Presentation

# Load presentation
prs = Presentation('output.pptx')
for i in range(0,len(prs.slides)):
    slide = prs.slides[i]
    print("Slide: ",i)
    # Iterate through all placeholders
    for shape in slide.placeholders:
        print(f"Placeholder idx: {shape.placeholder_format.idx}")
        print(f"Type: {shape.placeholder_format.type}")
        print(f"Name: {shape.name}")
        print("---")