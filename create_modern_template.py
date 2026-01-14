from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def create_modern_template():
    prs = Presentation()
    
    # 1. Title Slide (Modern Layout)
    # We use the default Title Slide but customize typical placement if we could,
    # but for a template, we just ensure the masters have the layouts we need.
    # By default, python-pptx creates a blank presentation which has standard layouts.
    # We will rename specific layouts to be clear for the AI.
    
    # Layout 0: Title Slide
    # Layout 1: Title and Content
    # Layout 8: Picture with Caption (This is usually good for images)
    
    # Let's ensure we have a good "Title + Picture + Text" layout.
    # We can't easily "create" a new master layout with pure python-pptx from scratch 
    # without manipulating XML, but we can utilize the existing ones effectively.
    
    # We'll save this as a "base" file that has been "touched" to look a bit modern 
    # (e.g. setting some theme colors is hard, but we can save it as a clean starting point).
    
    # Actually, a better approach is to ADD a slide with specific placeholders and then 
    # delete the slide, so the internal IDs are established? No.
    
    # We will just verify the layouts exist and maybe rename them if we could, 
    # but python-pptx doesn't let us rename layouts easily.
    
    # Instead, let's create a presentation that has content on it to guide the user,
    # or just save an empty one.
    
    # Wait, the best way to get a "Modern" template is to start with one.
    # Since I can't browse the web for a .pptx, I will create one that utilizes
    # the reliable layouts.
    
    title_slide_layout = prs.slide_layouts[0]
    title_and_content = prs.slide_layouts[1]
    two_content = prs.slide_layouts[3] # Good for Text + Image
    try:
        picture_caption = prs.slide_layouts[8] # Often Picture with Caption
    except:
        picture_caption = two_content
        
    # Let's verify what placeholders exist on "Two Content" (Layout 3)
    # Usually: Title, Content Placeholder 1 (Left), Content Placeholder 2 (Right)
    # Both content placeholders support images.
    
    # I'll save this file as 'templates_source/modern_template.pptx'
    # The app will list it automatically.
    
    prs.save('templates_source/modern_template.pptx')
    print("Created templates_source/modern_template.pptx")

if __name__ == "__main__":
    create_modern_template()
