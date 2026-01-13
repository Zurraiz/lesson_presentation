"""
Python-PPTX Template Explorer
This script demonstrates various slide layouts and features available in python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    """Create a PowerPoint presentation with various layout examples"""
    
    # Create presentation object
    prs = Presentation()
    
    # Get the slide layouts available in the default template
    print("Available slide layouts in default template:")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"  Layout {i}: {layout.name}")
    
    print("\nCreating presentation with various layouts...\n")
    
    # ========== Slide 1: Title Slide ==========
    slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Python-PPTX Template Explorer"
    subtitle.text = "Exploring Different Slide Layouts and Features\nGenerated with python-pptx"
    print("✓ Created Slide 1: Title Slide")
    
    # ========== Slide 2: Title and Content ==========
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Title and Content Layout"
    
    # Access the content placeholder
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "This is the main content area"
    
    # Add bullet points
    p = tf.add_paragraph()
    p.text = "First bullet point - placeholder text"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Second bullet point - more placeholder content"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Third bullet point - nested content"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "Fourth bullet point - back to main level"
    p.level = 0
    
    print("✓ Created Slide 2: Title and Content")
    
    # ========== Slide 3: Section Header ==========
    slide_layout = prs.slide_layouts[2]  # Section Header layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Section Header Layout"
    
    subtitle = slide.placeholders[1]
    subtitle.text = "This layout is great for dividing your presentation into sections"
    
    print("✓ Created Slide 3: Section Header")
    
    # ========== Slide 4: Two Content ==========
    slide_layout = prs.slide_layouts[3]  # Two Content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Two Content Layout"
    
    # Left content
    left_content = slide.placeholders[1]
    tf = left_content.text_frame
    tf.text = "Left Column Content"
    p = tf.add_paragraph()
    p.text = "Placeholder bullet 1"
    p = tf.add_paragraph()
    p.text = "Placeholder bullet 2"
    p = tf.add_paragraph()
    p.text = "Placeholder bullet 3"
    
    # Right content
    right_content = slide.placeholders[2]
    tf = right_content.text_frame
    tf.text = "Right Column Content"
    p = tf.add_paragraph()
    p.text = "Placeholder item A"
    p = tf.add_paragraph()
    p.text = "Placeholder item B"
    p = tf.add_paragraph()
    p.text = "Placeholder item C"
    
    print("✓ Created Slide 4: Two Content")
    
    # ========== Slide 5: Comparison ==========
    slide_layout = prs.slide_layouts[4]  # Comparison layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Comparison Layout"
    
    # This layout varies by template, but typically has comparison structure
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            shape.text = "Option A\n• Feature 1\n• Feature 2\n• Feature 3"
        elif shape.placeholder_format.idx == 2:
            shape.text = "Option B\n• Feature X\n• Feature Y\n• Feature Z"
    
    print("✓ Created Slide 5: Comparison")
    
    # ========== Slide 6: Title Only ==========
    slide_layout = prs.slide_layouts[5]  # Title Only layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Title Only Layout"
    
    # Add custom text box in the content area
    left = Inches(1.5)
    top = Inches(2.5)
    width = Inches(7)
    height = Inches(3)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.text = "This layout gives you a blank canvas below the title.\nYou can add custom shapes, text boxes, and images anywhere you want."
    
    # Format the text
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER
    
    print("✓ Created Slide 6: Title Only")
    
    # ========== Slide 7: Blank Layout ==========
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title manually
    left = Inches(1)
    top = Inches(0.5)
    width = Inches(8)
    height = Inches(1)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    tf.text = "Blank Layout - Completely Customizable"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    # Add some content
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(3.5)
    
    content_box = slide.shapes.add_textbox(left, top, width, height)
    tf = content_box.text_frame
    tf.text = "This is a blank layout where you have full control.\n\n"
    tf.text += "You can add:\n"
    tf.text += "• Text boxes (like this one)\n"
    tf.text += "• Shapes (rectangles, circles, etc.)\n"
    tf.text += "• Images (placeholder images)\n"
    tf.text += "• Charts and tables\n"
    tf.text += "• And position them anywhere!"
    
    print("✓ Created Slide 7: Blank Layout")
    
    # ========== Slide 8: Custom Shapes Demo ==========
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Custom Shapes and Formatting"
    
    # Add rectangles with different colors
    from pptx.enum.shapes import MSO_SHAPE
    
    left = Inches(1)
    top = Inches(2)
    width = Inches(2)
    height = Inches(1.5)
    
    # Blue rectangle
    shape1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RGBColor(0, 112, 192)
    shape1.text_frame.text = "Placeholder\nBox 1"
    shape1.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    shape1.text_frame.paragraphs[0].font.size = Pt(18)
    shape1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Green rectangle
    shape2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left + Inches(2.5), top, width, height
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(0, 176, 80)
    shape2.text_frame.text = "Placeholder\nBox 2"
    shape2.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    shape2.text_frame.paragraphs[0].font.size = Pt(18)
    shape2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Orange rectangle
    shape3 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left + Inches(5), top, width, height
    )
    shape3.fill.solid()
    shape3.fill.fore_color.rgb = RGBColor(255, 127, 39)
    shape3.text_frame.text = "Placeholder\nBox 3"
    shape3.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    shape3.text_frame.paragraphs[0].font.size = Pt(18)
    shape3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add a rounded rectangle
    shape4 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(4), Inches(5), Inches(1)
    )
    shape4.fill.solid()
    shape4.fill.fore_color.rgb = RGBColor(128, 0, 128)
    shape4.text_frame.text = "Rounded Rectangle - Placeholder for important note"
    shape4.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    shape4.text_frame.paragraphs[0].font.size = Pt(16)
    shape4.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    print("✓ Created Slide 8: Custom Shapes")
    
    # ========== Slide 9: Picture Placeholder Demo ==========
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Picture Placeholders"
    
    # Add placeholder text boxes where images would go
    img_positions = [
        (Inches(1), Inches(2), Inches(3), Inches(2.5)),
        (Inches(5.5), Inches(2), Inches(3), Inches(2.5)),
    ]
    
    for i, (left, top, width, height) in enumerate(img_positions, 1):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(220, 220, 220)
        shape.line.color.rgb = RGBColor(150, 150, 150)
        shape.text_frame.text = f"[Image Placeholder {i}]\n\nYou can insert images here\nusing add_picture()"
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        shape.text_frame.paragraphs[0].font.size = Pt(14)
        shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
    
    print("✓ Created Slide 9: Picture Placeholders")
    
    # ========== Slide 10: Summary ==========
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Python-PPTX Capabilities Summary"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "What we've explored:"
    
    capabilities = [
        "Multiple built-in slide layouts (Title, Content, Section Header, etc.)",
        "Text formatting with fonts, sizes, colors, and alignment",
        "Bullet points and nested lists",
        "Custom text boxes and positioning",
        "Shapes (rectangles, rounded rectangles, and more)",
        "Color customization (fill colors, text colors)",
        "Placeholders for images and content",
        "Full layout customization on blank slides"
    ]
    
    for capability in capabilities:
        p = tf.add_paragraph()
        p.text = capability
        p.level = 0
    
    print("✓ Created Slide 10: Summary")
    
    # Save the presentation
    filename = "python_pptx_template_demo.pptx"
    prs.save(filename)
    print(f"\n✅ Presentation saved as '{filename}'")
    print(f"   Total slides created: {len(prs.slides)}")
    
    return filename

if __name__ == "__main__":
    print("=" * 60)
    print("Python-PPTX Template Explorer")
    print("=" * 60)
    print()
    
    filename = create_presentation()
    
    print()
    print("=" * 60)
    print(f"✅ Successfully created: {filename}")
    print("Open this file to see all the different layouts and features!")
    print("=" * 60)
