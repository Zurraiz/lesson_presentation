"""
Python-PPTX Template Manager - Refactored
Clean API for loading templates, inspecting placeholders, and filling them with content
"""

from pptx import Presentation
from pptx.util import Pt
import os
from typing import Dict, List, Any, Optional


class TemplateManager:
    """Manages PowerPoint template operations"""
    
    def __init__(self, template_path: str):
        """
        Initialize with a template file
        
        Args:
            template_path: Path to the .pptx template file
        """
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template not found: {template_path}")
        
        self.template_path = template_path
        self.presentation = Presentation(template_path)
    
    def get_placeholder_info(self, slide_index: int = None) -> Dict[str, Any]:
        """
        Get detailed information about placeholders in the template
        
        Args:
            slide_index: If provided, get info for specific slide. Otherwise, get template layout info
            
        Returns:
            Dictionary containing placeholder information
        """
        if slide_index is not None:
            # Get info for a specific slide
            if slide_index >= len(self.presentation.slides):
                return {"error": f"Slide {slide_index} does not exist"}
            
            slide = self.presentation.slides[slide_index]
            return self._extract_slide_placeholder_info(slide, slide_index)
        else:
            # Get info for all layouts in template
            return self._extract_layout_info()
    
    def _extract_layout_info(self) -> Dict[str, Any]:
        """Extract placeholder information from all layouts"""
        layout_info = {
            "template_path": self.template_path,
            "total_layouts": len(self.presentation.slide_layouts),
            "layouts": []
        }
        
        for idx, layout in enumerate(self.presentation.slide_layouts):
            layout_data = {
                "index": idx,
                "name": layout.name,
                "placeholders": []
            }
            
            for placeholder in layout.placeholders:
                ph_info = {
                    "idx": placeholder.placeholder_format.idx,
                    "name": placeholder.name,
                    "type": str(placeholder.placeholder_format.type),
                    "has_text_frame": placeholder.has_text_frame,
                }
                
                if placeholder.has_text_frame and placeholder.text:
                    ph_info["default_text"] = placeholder.text
                
                layout_data["placeholders"].append(ph_info)
            
            layout_info["layouts"].append(layout_data)
        
        return layout_info
    
    def _extract_slide_placeholder_info(self, slide, slide_index: int) -> Dict[str, Any]:
        """Extract placeholder information from a specific slide"""
        slide_info = {
            "slide_index": slide_index,
            "placeholders": []
        }
        
        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_info = {
                    "idx": shape.placeholder_format.idx,
                    "name": shape.name,
                    "type": str(shape.placeholder_format.type),
                    "has_text_frame": shape.has_text_frame,
                }
                
                if shape.has_text_frame:
                    ph_info["current_text"] = shape.text
                
                slide_info["placeholders"].append(ph_info)
        
        return slide_info
    
    def print_placeholder_info(self, slide_index: int = None):
        """Print placeholder information in a readable format"""
        info = self.get_placeholder_info(slide_index)
        
        if slide_index is None:
            # Print layout information
            print(f"\n{'='*70}")
            print(f"TEMPLATE: {info['template_path']}")
            print(f"{'='*70}\n")
            print(f"Total Layouts: {info['total_layouts']}\n")
            
            for layout in info['layouts']:
                print(f"Layout {layout['index']}: {layout['name']}")
                print("-" * 70)
                if not layout['placeholders']:
                    print("  No placeholders\n")
                    continue
                
                for ph in layout['placeholders']:
                    print(f"  Placeholder {ph['idx']}: {ph['name']}")
                    print(f"    Type: {ph['type']}")
                    print(f"    Has Text Frame: {ph['has_text_frame']}")
                    if 'default_text' in ph:
                        print(f"    Default: '{ph['default_text'][:50]}...'")
                    print()
        else:
            # Print slide information
            print(f"\nSlide {info['slide_index']} Placeholders:")
            print("-" * 70)
            for ph in info['placeholders']:
                print(f"  Placeholder {ph['idx']}: {ph['name']}")
                print(f"    Type: {ph['type']}")
                if 'current_text' in ph:
                    print(f"    Text: '{ph['current_text'][:50]}...'")
                print()
    
    def add_slide(self, layout_index: int, placeholder_data: Dict[int, str]) -> int:
        """
        Add a new slide with specified content
        
        Args:
            layout_index: Index of the layout to use
            placeholder_data: Dictionary mapping placeholder index to content
            
        Returns:
            Index of the newly created slide
        """
        if layout_index >= len(self.presentation.slide_layouts):
            raise ValueError(f"Layout {layout_index} does not exist")
        
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[layout_index])
        
        # Fill placeholders
        for ph_idx, content in placeholder_data.items():
            for shape in slide.shapes:
                if shape.is_placeholder and shape.placeholder_format.idx == ph_idx:
                    if shape.has_text_frame:
                        shape.text = content
                    break
        
        return len(self.presentation.slides) - 1
    
    def fill_placeholder(self, slide_index: int, placeholder_index: int, content: str) -> bool:
        """
        Fill a specific placeholder in a specific slide
        
        Args:
            slide_index: Index of the slide
            placeholder_index: Index of the placeholder
            content: Content to fill
            
        Returns:
            True if successful, False otherwise
        """
        if slide_index >= len(self.presentation.slides):
            print(f"❌ Slide {slide_index} does not exist")
            return False
        
        slide = self.presentation.slides[slide_index]
        
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == placeholder_index:
                if shape.has_text_frame:
                    shape.text = content
                    return True
                else:
                    print(f"❌ Placeholder {placeholder_index} has no text frame")
                    return False
        
        print(f"❌ Placeholder {placeholder_index} not found in slide {slide_index}")
        return False
    
    def save(self, output_path: str):
        """Save the presentation to a file"""
        self.presentation.save(output_path)
        print(f"✅ Saved presentation to: {output_path}")


def generate_dummy_content(layout_name: str, placeholder_info: Dict) -> str:
    """
    Generate dummy content based on placeholder type and position
    
    Args:
        layout_name: Name of the layout
        placeholder_info: Placeholder information dictionary
        
    Returns:
        Dummy content string
    """
    ph_name = placeholder_info['name'].lower()
    ph_type = placeholder_info['type']
    
    # Generate content based on placeholder type
    if 'title' in ph_name:
        return f"Sample Title for {layout_name}"
    elif 'subtitle' in ph_name:
        return f"Subtitle: Auto-generated content for demonstration"
    elif 'content' in ph_name or 'body' in ph_name or 'text' in ph_name:
        return (
            "Automatically generated content\n\n"
            "• Bullet point 1 - Sample text\n"
            "• Bullet point 2 - More sample text\n"
            "• Bullet point 3 - Additional content\n"
            "• Bullet point 4 - Final point"
        )
    elif 'picture' in ph_name or 'image' in ph_name:
        return "[Image placeholder - would insert image here]"
    elif 'date' in ph_name:
        return "January 13, 2026"
    elif 'footer' in ph_name:
        return "Auto-generated Presentation"
    elif 'slide number' in ph_name or 'number' in ph_name:
        return "1"
    else:
        return f"Sample content for {placeholder_info['name']}"


def create_sample_presentation_with_dummy_data(template_path: str, output_path: str, num_slides: int = 5):
    """
    Example function that creates a presentation with automatically generated dummy data
    
    Args:
        template_path: Path to the template file
        output_path: Path where to save the filled presentation
        num_slides: Number of slides to create
    """
    print(f"\n{'='*70}")
    print("CREATING PRESENTATION WITH DUMMY DATA")
    print(f"{'='*70}\n")
    
    # Create template manager
    tm = TemplateManager(template_path)
    
    # Get layout information
    layout_info = tm.get_placeholder_info()
    
    # Select layouts to use (cycle through available layouts)
    available_layouts = [l for l in layout_info['layouts'] if l['placeholders']]
    
    if not available_layouts:
        print("❌ No layouts with placeholders found")
        return
    
    # Add slides with dummy content
    for i in range(num_slides):
        layout = available_layouts[i % len(available_layouts)]
        layout_idx = layout['index']
        layout_name = layout['name']
        
        print(f"Adding slide {i+1} using layout '{layout_name}'...")
        
        # Generate dummy content for each placeholder
        placeholder_data = {}
        for ph in layout['placeholders']:
            if ph['has_text_frame']:
                content = generate_dummy_content(layout_name, ph)
                placeholder_data[ph['idx']] = content
                print(f"  Placeholder {ph['idx']} ({ph['name']}): {content[:40]}...")
        
        # Add the slide
        tm.add_slide(layout_idx, placeholder_data)
        print()
    
    # Save the presentation
    tm.save(output_path)
    print(f"\n✅ Created {num_slides} slides with dummy data")


def create_sample_template():
    """Create a sample template for testing"""
    print("Creating sample template...")
    prs = Presentation()
    
    # Add various layout slides
    prs.slides.add_slide(prs.slide_layouts[0])  # Title
    prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    prs.slides.add_slide(prs.slide_layouts[3])  # Two Content
    prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    
    template_name = "sample_template.pptx"
    prs.save(template_name)
    print(f"✓ Created: {template_name}\n")
    return template_name


if __name__ == "__main__":
    print("\n" + "="*70)
    print("PYTHON-PPTX TEMPLATE MANAGER - REFACTORED")
    print("="*70)
    
    # 1. Create a sample template
    template_path = create_sample_template()
    
    # 2. Inspect the template
    print("\n--- INSPECTING TEMPLATE ---")
    tm = TemplateManager(template_path)
    tm.print_placeholder_info()
    
    # 3. Create presentation with dummy data (automatic content generation)
    output_path = "auto_generated_presentation.pptx"
    create_sample_presentation_with_dummy_data(template_path, output_path, num_slides=5)
    
    # 4. Show example of manual control
    print("\n--- EXAMPLE: MANUAL CONTROL ---")
    print("You can also manually control content:\n")
    
    tm2 = TemplateManager(template_path)
    
    # Add a slide manually
    slide_idx = tm2.add_slide(layout_index=1, placeholder_data={
        0: "Custom Title",
        1: "Custom content added manually"
    })
    print(f"Added custom slide at index {slide_idx}")
    
    # Modify existing slide
    tm2.fill_placeholder(slide_index=0, placeholder_index=0, content="Modified Title")
    print("Modified existing slide")
    
    tm2.save("manually_created_presentation.pptx")
    
    print("\n" + "="*70)
    print("COMPLETE!")
    print("="*70)
    print("\nGenerated files:")
    print(f"  • {template_path}")
    print(f"  • {output_path} (auto-generated with dummy data)")
    print(f"  • manually_created_presentation.pptx")
    print("\nAPI Usage:")
    print("  tm = TemplateManager('template.pptx')")
    print("  tm.get_placeholder_info()  # Get placeholder details")
    print("  tm.add_slide(layout_idx, {placeholder_idx: 'content'})")
    print("  tm.fill_placeholder(slide_idx, placeholder_idx, 'content')")
    print("  tm.save('output.pptx')")
    print("="*70 + "\n")
